from pptx import Presentation
from pptx.util import Inches, Pt

slides_content = [
    ("NexusCare: Full-Stack Digital Ecosystem", [
        "Smart community services for residents, security, medical, admin",
        "Stack: React (Vite) + Flask API + MySQL",
        "Features: Auth, role-based access, Complaints CRUD, audit logs",
    ]),
    ("Problem & Objectives", [
        "Challenge: Multi-role, secure, scalable community platform",
        "Objectives: Reliable auth, protected routes, persistent storage, modularity",
        "Outcomes: Production-ready patterns, clean APIs, responsive UI",
    ]),
    ("System Architecture", [
        "SPA frontend → Axios → REST API /api/* → MySQL",
        "Sessions via Flask-Login; CORS via Flask-CORS; rotating file logs",
        "Vite dev proxy: /api → http://localhost:5000",
    ]),
    ("Backend Design (Flask)", [
        "app.py: factory, CORS, SQLAlchemy, LoginManager, error handlers, health",
        "Blueprints: routes/auth.py, routes/complaints.py",
        "Config: config.py; Extensions: extensions.py (db, login)",
        "Logging: instance/logs/app.log",
    ]),
    ("Database & Models (MySQL + SQLAlchemy)", [
        "Tables: users, complaints, audit_logs",
        "Models: User, Complaint, AuditLog; enums for Role and ComplaintStatus",
        "Relations: User 1–N Complaint; audit entries per action",
        "Init: schema.sql and init_db.py (DB/user/grants)",
    ]),
    ("Security & Authorization", [
        "Authentication: Flask-Login + password hashing",
        "Session endpoints: POST /auth/login, POST /auth/logout, GET /auth/me",
        "Role checks: Admin/security elevated; resident/medical own records",
        "Errors: 401/403/404/500 JSON",
    ]),
    ("Frontend (React + Vite)", [
        "Routing: App.jsx with react-router-dom and ProtectedRoute",
        "State: AuthContext manages session via /auth/me",
        "Pages: Login, Register, Dashboard, Complaints, Unauthorized, NotFound",
        "UI: NavBar, reusable components, responsive styles.css",
    ]),
    ("Complaints Module (CRUD)", [
        "Create: POST /api/complaints; Read: GET /api/complaints?status=",
        "Update: PUT /api/complaints/:id; Delete (soft): DELETE /api/complaints/:id",
        "Audit logging on create/update/delete",
        "UI: filter by status, inline status updates, soft delete",
    ]),
    ("Deployment, Testing & Roadmap", [
        "Run: python init_db.py → python app.py; npm install → npm run dev",
        "API tests: Postman or Invoke-RestMethod",
        "Next: modules (visitors, announcements, appointments), pagination, validation, Docker, CI/CD",
    ]),
]

prs = Presentation()
# Use Title and Content layout
layout = prs.slide_layouts[1]

for title, bullets in slides_content:
    slide = prs.slides.add_slide(layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    tf = body_placeholder.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    # first bullet as paragraph
    p = tf.paragraphs[0]
    p.text = bullets[0]
    p.font.size = Pt(20)

    for bullet in bullets[1:]:
        para = tf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.font.size = Pt(18)

prs.save("NexusCare_Deck.pptx")
print("Saved NexusCare_Deck.pptx")
